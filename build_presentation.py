"""Генерирует presentation.pptx с защитой проекта по кейсу 09 и README."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "presentation.pptx"

BG = RGBColor(0x0F, 0x11, 0x15)
FG = RGBColor(0xE8, 0xEA, 0xED)
MUTED = RGBColor(0x9A, 0xA0, 0xA6)
ACCENT = RGBColor(0x7A, 0xA2, 0xFF)
CARD = RGBColor(0x18, 0x1B, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW, SH = prs.slide_width, prs.slide_height


def paint_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, left, top, width, height, text, *, size=18, color=FG, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Calibri"
    return tb


def add_bullets(slide, left, top, width, height, items, *, size=18, color=FG):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_header(slide, title):
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.1), Inches(0.5))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    add_text(slide, Inches(0.85), Inches(0.5), Inches(12), Inches(0.7), title, size=28, bold=True)


def card(slide, left, top, width, height, title, body_lines, *, body_size=15):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.06
    box.line.color.rgb = RGBColor(0x24, 0x28, 0x33)
    box.line.width = Pt(0.75)
    box.fill.solid()
    box.fill.fore_color.rgb = CARD
    add_text(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.4), Inches(0.45),
             title, size=17, bold=True, color=ACCENT)
    tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.65),
                                  width - Inches(0.4), height - Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(body_size)
        run.font.color.rgb = FG
        run.font.name = "Calibri"


def new_slide(title=None):
    s = prs.slides.add_slide(BLANK)
    paint_bg(s)
    if title:
        add_header(s, title)
    return s


# ---------- Slide 1: Title ----------
s = new_slide()
add_text(s, Inches(0.8), Inches(2.2), Inches(11.8), Inches(1.2),
         "Система семантического поиска\nпо базе знаний компании",
         size=44, bold=True, align=PP_ALIGN.LEFT)
add_text(s, Inches(0.8), Inches(3.8), Inches(11.8), Inches(0.5),
         "MVP • Кейс 09 • Прикладной искусственный интеллект", size=20, color=ACCENT)
add_text(s, Inches(0.8), Inches(6.3), Inches(11.8), Inches(0.5),
         "Направление 09.03.03 — Прикладная информатика, 2-й семестр", size=14, color=MUTED)
add_text(s, Inches(0.8), Inches(6.7), Inches(11.8), Inches(0.5),
         "GitHub: github.com/ustyusty/semantic-search", size=14, color=MUTED)

# ---------- Slide 2: Проблема ----------
s = new_slide("Проблема и актуальность")
add_text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2),
         "В любой компании накапливаются сотни регламентов и инструкций.\n"
         "Классический поиск (Ctrl+F) не работает, если сотрудник не знает\n"
         "точного названия документа или термина.",
         size=22)
add_text(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.6),
         "Семантический поиск решает это: ищет по смыслу, а не по словам.", size=20, color=ACCENT, bold=True)
card(s, Inches(0.8), Inches(4.8), Inches(5.8), Inches(1.9),
     "Запрос пользователя",
     ['«Хочу отдохнуть»', '«Что делать если сломался принтер»', '«Где поесть»'])
card(s, Inches(7.0), Inches(4.8), Inches(5.8), Inches(1.9),
     "Найденный документ",
     ['«Инструкция: Как оформить отпуск»', '«Правила эксплуатации офисной техники»', '«Расположение столовой»'])

# ---------- Slide 3: Цель и результат ----------
s = new_slide("Цель и конкретный результат")
add_text(s, Inches(0.8), Inches(1.45), Inches(12), Inches(1),
         "Работающий прототип веб-сервиса для поиска документов\nпо смысловой близости к запросу.",
         size=22)
add_bullets(s, Inches(0.9), Inches(3.5), Inches(12), Inches(3.5), [
    "Администратор — загружает документы (текст или PDF) в базу знаний",
    "Пользователь — вводит запрос в свободной форме",
    "Система — возвращает топ-3 наиболее релевантных документа",
    "Без обучения нейросетей с нуля — используется готовая предобученная модель",
], size=20)

# ---------- Slide 4: Стек технологий ----------
s = new_slide("Стек технологий")
cards = [
    ("Backend", ["Python 3.10+", "FastAPI + Uvicorn", "asyncpg (async Postgres)", "Pydantic для валидации"]),
    ("NLP / ML", ["sentence-transformers", "paraphrase-multilingual-MiniLM-L12-v2", "NumPy для матричных операций", "pypdf для извлечения текста из PDF"]),
    ("Инфраструктура", ["PostgreSQL 15", "Docker Compose", "HTTP Basic Auth", "Git / GitHub"]),
]
for i, (t, items) in enumerate(cards):
    card(s, Inches(0.7 + i * 4.2), Inches(1.7), Inches(4.0), Inches(5.4), t, items, body_size=16)

# ---------- Slide 5: Архитектура ----------
s = new_slide("Архитектура")
# Три блока: клиент → FastAPI → Postgres
blocks = [
    ("Клиент", "Браузер / curl\nHTML + JS (vanilla)", Inches(0.8), Inches(2.2)),
    ("Сервис", "FastAPI (lifespan)\nРоуты + Basic Auth\nВ памяти: NumPy-индекс", Inches(5.1), Inches(2.2)),
    ("Хранилище", "PostgreSQL\nТаблица documents\nembedding DOUBLE[]", Inches(9.4), Inches(2.2)),
]
for title, body, left, top in blocks:
    card(s, left, top, Inches(3.3), Inches(2.5), title, body.split("\n"))

add_text(s, Inches(0.8), Inches(5.3), Inches(12), Inches(0.5),
         "Поток поиска:", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(0.9), Inches(5.8), Inches(12), Inches(1.5), [
    "Клиент → GET /api/search?q=… → FastAPI",
    "Эмбеддер превращает запрос в вектор (384 dim, normalized)",
    "NumPy: matrix @ query_vec → сортировка → топ-K → JSON",
], size=16)

# ---------- Slide 6: Как работает поиск ----------
s = new_slide("Как работает семантический поиск")
add_bullets(s, Inches(0.8), Inches(1.6), Inches(12), Inches(5), [
    "Каждый документ превращается в вектор фиксированной длины (embedding, 384 числа)",
    "Модель paraphrase-multilingual-MiniLM-L12-v2 размещает похожие по смыслу тексты рядом",
    "Близость двух векторов — косинусное сходство: cos(a,b) = (a·b) / (|a|·|b|)",
    "Эмбеддинги нормализованы — косинус = скалярное произведение, одно матричное умножение",
    "Поиск: один вызов matrix @ query_vec находит расстояния до всех документов за миллисекунды",
    "Результат: отсортированный топ-K с оценкой от 0 до 1",
], size=18)

# ---------- Slide 7: Схема БД ----------
s = new_slide("Схема базы данных")
sql = (
    "CREATE TABLE documents (\n"
    "    id          SERIAL PRIMARY KEY,\n"
    "    title       TEXT NOT NULL,\n"
    "    content     TEXT NOT NULL,\n"
    "    block       TEXT,\n"
    "    embedding   DOUBLE PRECISION[],\n"
    "    created_at  TIMESTAMP DEFAULT NOW()\n"
    ");"
)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(7.5), Inches(3.8))
box.adjustments[0] = 0.04
box.fill.solid(); box.fill.fore_color.rgb = CARD
box.line.color.rgb = RGBColor(0x24, 0x28, 0x33)
tb = s.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(7.2), Inches(3.5))
tf = tb.text_frame; tf.word_wrap = True
for i, line in enumerate(sql.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run(); run.text = line
    run.font.name = "Consolas"; run.font.size = Pt(16); run.font.color.rgb = RGBColor(0xD0, 0xD4, 0xDB)

add_text(s, Inches(8.6), Inches(1.7), Inches(4.2), Inches(0.5),
         "Особенности", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(8.6), Inches(2.2), Inches(4.4), Inches(4), [
    "Embedding в виде массива float",
    "При старте всё грузится в NumPy-матрицу",
    "Блок — логическая группировка (HR, IT …)",
    "Индекс обновляется при каждой вставке",
], size=15)

# ---------- Slide 8: API ----------
s = new_slide("API: публичное и админское")
card(s, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.4),
     "Публичные эндпоинты",
     [
         "GET /  — страница поиска (HTML)",
         "GET /api/search?q=…&k=3",
         "GET /api/documents/count",
         "GET /docs  — Swagger UI",
         "",
         "Не требуют авторизации.",
         "Любой пользователь может искать.",
     ])
card(s, Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.4),
     "Админские эндпоинты (Basic Auth)",
     [
         "GET  /admin  — форма загрузки",
         "POST /api/documents  — JSON {title, content, block}",
         "POST /api/documents/upload  — PDF, multipart",
         "",
         "Пароль в .env (ADMIN_USER, ADMIN_PASSWORD).",
         "Без него — 401 Unauthorized.",
     ])

# ---------- Slide 9: UI поиска ----------
s = new_slide("Интерфейс пользователя")
add_bullets(s, Inches(0.8), Inches(1.6), Inches(12), Inches(4), [
    "Минималистичная страница с поисковой строкой",
    "Запрос на естественном языке — «хочу отдохнуть», «сломался принтер»",
    "Топ-3 документа в виде карточек: заголовок, раздел, фрагмент текста, score",
    "Показ общего количества документов в базе",
    "Счётчик релевантности — для прозрачности качества поиска",
], size=18)
add_text(s, Inches(0.8), Inches(5.3), Inches(12), Inches(0.5),
         "URL: http://<host>:8000/", size=16, color=ACCENT)

# ---------- Slide 10: Админка ----------
s = new_slide("Админка")
add_bullets(s, Inches(0.8), Inches(1.6), Inches(12), Inches(4.5), [
    "Вход по HTTP Basic Auth (прямо через браузерный диалог)",
    "Две вкладки загрузки:",
    "    — «Текст»: поля Заголовок, Раздел, Текст документа",
    "    — «PDF-файл»: загрузка PDF (до 15 МБ), парсинг через pypdf",
    "При добавлении эмбеддинг считается автоматически",
    "Индекс обновляется — документ сразу доступен для поиска",
    "PDF без извлекаемого текста (скан без OCR) — вежливо отклоняется",
], size=17)
add_text(s, Inches(0.8), Inches(6.0), Inches(12), Inches(0.5),
         "URL: http://<host>:8000/admin", size=16, color=ACCENT)

# ---------- Slide 11: Безопасность ----------
s = new_slide("Безопасность")
add_bullets(s, Inches(0.8), Inches(1.6), Inches(12), Inches(4.5), [
    "Роли по требованию кейса: только администратор добавляет документы",
    "FastAPI-dependency require_admin подключена к POST-роутам и к /admin",
    "Сравнение пароля через secrets.compare_digest (устойчиво к timing-атакам)",
    "Шаблон admin.html вынесен из static/ — нельзя обойти через прямой URL",
    "Учётные данные — в .env, не в коде и не в git",
    "Обычный пользователь видит только страницу поиска и /docs",
], size=18)

# ---------- Slide 12: Быстрый старт ----------
s = new_slide("Быстрый старт (из README)")
steps = [
    ("1", "cp .env.example .env  &&  docker compose up -d"),
    ("2", "python -m venv .venv && source .venv/bin/activate"),
    ("3", "pip install -r requirements.txt"),
    ("4", "uvicorn app.main:app --host 0.0.0.0 --port 8000"),
    ("5", "Открыть http://localhost:8000  —  поиск"),
    ("6", "Открыть http://localhost:8000/admin  —  загрузка документов"),
]
top = Inches(1.6)
for num, cmd in steps:
    num_box = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), top, Inches(0.55), Inches(0.55))
    num_box.fill.solid(); num_box.fill.fore_color.rgb = ACCENT
    num_box.line.fill.background()
    tf = num_box.text_frame; tf.margin_top = tf.margin_bottom = Pt(0)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = num
    r.font.bold = True; r.font.size = Pt(18); r.font.color.rgb = RGBColor(0x0B, 0x0D, 0x12)
    add_text(s, Inches(1.6), top + Inches(0.05), Inches(11.5), Inches(0.5), cmd, size=16)
    top += Inches(0.75)

# ---------- Slide 13: Пример работы ----------
s = new_slide("Пример работы поиска")
add_text(s, Inches(0.8), Inches(1.5), Inches(12), Inches(0.5),
         'Запрос: «хочу отдохнуть»  →  найдено:', size=20, bold=True)
card(s, Inches(0.8), Inches(2.3), Inches(12), Inches(1.4),
     "1. Инструкция: Как оформить отпуск   (score 0.29)",
     ["Каждый сотрудник имеет право на 28 календарных дней ежегодного оплачиваемого отпуска.",
      "График отпусков утверждается в корпоративной ERP-системе до 15 декабря…"])
add_text(s, Inches(0.8), Inches(3.9), Inches(12), Inches(0.5),
         'Запрос: «где поесть»  →  найдено:', size=20, bold=True)
card(s, Inches(0.8), Inches(4.7), Inches(12), Inches(1.4),
     "1. Правила пользования кофе-пойнтом   (score 0.41)",
     ["Просьба мыть за собой чашки и подписывать контейнеры в холодильнике.",
      "Слова «поесть» в тексте нет — нашли по смыслу."])

# ---------- Slide 14: Milestones ----------
s = new_slide("Покрытие milestones кейса")
milestones = [
    ("Показ 1 — данные и стек", "Выбор sentence-transformers, документы в JSON"),
    ("Показ 2 — админка", "POST /api/documents, сохранение в PostgreSQL"),
    ("Показ 3 — мозги поиска", "Эмбеддинги + косинусное сходство, топ-3 в Swagger"),
    ("Показ 4 — MVP", "Веб-интерфейс, связанный с API поиска"),
    ("Финал — тестирование", "README, примеры сложных запросов, презентация"),
]
top = Inches(1.6)
for title, body in milestones:
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(0.95))
    box.adjustments[0] = 0.25
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.color.rgb = RGBColor(0x24, 0x28, 0x33)
    add_text(s, Inches(1.0), top + Inches(0.1), Inches(4.5), Inches(0.4), title, size=16, bold=True, color=ACCENT)
    add_text(s, Inches(1.0), top + Inches(0.5), Inches(11.3), Inches(0.4), body, size=14)
    top += Inches(1.05)

# ---------- Slide 15: Критерии оценки ----------
s = new_slide("Критерии оценки кейса")
add_bullets(s, Inches(0.8), Inches(1.5), Inches(12), Inches(5), [
    "Качество поиска (20) — находит по смыслу, не только по совпадению слов",
    "Простота развертывания (10) — docker compose up -d + uvicorn, одна команда",
    "Интерфейс (10) — минималистичный, две страницы (поиск и админка)",
    "Презентация (10) — этот файл",
    "Ответы на вопросы (10) — понимание эмбеддингов и косинусного сходства",
], size=19)
add_text(s, Inches(0.8), Inches(5.8), Inches(12), Inches(0.5),
         "Ожидаемая оценка: Отлично — поиск работает семантически, UI чистый, код структурирован.",
         size=16, color=ACCENT)

# ---------- Slide 16: Итог ----------
s = new_slide("Итог")
add_text(s, Inches(0.8), Inches(2.0), Inches(12), Inches(1),
         "MVP системы семантического поиска готов.", size=30, bold=True)
add_bullets(s, Inches(0.8), Inches(3.2), Inches(12), Inches(3), [
    "Развёрнут на сервере, доступен снаружи",
    "Поддержка загрузки PDF и чистого текста",
    "Разделение ролей: пользователь / администратор",
    "Открытый исходный код на GitHub",
], size=18)
add_text(s, Inches(0.8), Inches(6.3), Inches(12), Inches(0.5),
         "Спасибо за внимание!", size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"OK: {OUT}  ({OUT.stat().st_size // 1024} KB)")
